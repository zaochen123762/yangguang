import os
import streamlit as st
from zhipuai import ZhipuAI
import PyPDF2
import chromadb
import requests
from langchain_community.embeddings import HuggingFaceEmbeddings
import docx
from pptx import Presentation
import openpyxl
API_KEY = "db4b7ce1e5384d96b348a877bf0c8f60.T300s1a8zdliQxq2"
client = ZhipuAI(api_key=API_KEY)
@st.cache_resource
def init_chroma():
    chroma_client = chromadb.PersistentClient(path="./chroma_db")
    embedding_fn = HuggingFaceEmbeddings(
        model_name="sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2"
    )
    collection = chroma_client.get_or_create_collection(name="documents")
    return chroma_client, collection, embedding_fn
def extract_text_from_file(uploaded_file):
    """根据文件类型解析文本内容，支持 PDF / Word / PPT / Excel / TXT"""
    file_name = uploaded_file.name
    text = ""
    if file_name.endswith('.pdf'):
        reader = PyPDF2.PdfReader(uploaded_file)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    elif file_name.endswith('.docx'):
        doc = docx.Document(uploaded_file)
        for para in doc.paragraphs:
            if para.text:
                text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                row_text = ""
                for cell in row.cells:
                    if cell.text:
                        row_text += cell.text + " "
                if row_text.strip():
                    text += row_text + "\n"
    elif file_name.endswith('.pptx'):
        prs = Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    elif file_name.endswith('.xlsx'):
        wb = openpyxl.load_workbook(uploaded_file, data_only=True)
        for sheet in wb.worksheets:
            text += f"=== {sheet.title} ===\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = ""
                for cell in row:
                    if cell is not None:
                        row_text += str(cell) + " "
                if row_text.strip():
                    text += row_text + "\n"
    elif file_name.endswith('.txt'):
        text = uploaded_file.read().decode('utf-8')
    else:
        raise ValueError(f"不支持的文件格式: {file_name}")
    return text
def split_text(text, chunk_size=500):
    chunks = []
    paragraphs = text.split("\n\n")
    current = ""
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if len(current) + len(para) < chunk_size:
            current += para + "\n"
        else:
            if current.strip():
                chunks.append(current.strip())
            current = para + "\n"
    if current.strip():
        chunks.append(current.strip())
    return chunks
def upload_document(file, collection, embedding_fn):
    try:
        text = extract_text_from_file(file)
        if not text.strip():
            return 0, "文件内容为空或无法解析"
        chunks = split_text(text)
        if not chunks:
            return 0, "未提取到有效文本内容"
        doc_id = file.name
        for ext in ['.pdf', '.docx', '.pptx', '.xlsx', '.txt']:
            doc_id = doc_id.replace(ext, '')
        for i, chunk in enumerate(chunks):
            embedding = embedding_fn.embed_query(chunk)
            chunk_id = f"{doc_id}_{i}"
            collection.add(
                documents=[chunk],
                embeddings=[embedding],
                ids=[chunk_id],
                metadatas=[{"source": file.name, "chunk": i}]
            )
        return len(chunks), "成功"
    except Exception as e:
        return 0, f"解析失败: {str(e)}"
def search_documents(query, collection, embedding_fn, top_k=3):
    query_embedding = embedding_fn.embed_query(query)
    results = collection.query(
        query_embeddings=[query_embedding],
        n_results=top_k
    )
    if results['documents'] and len(results['documents'][0]) > 0:
        contexts = []
        for i, doc in enumerate(results['documents'][0]):
            source = results['metadatas'][0][i]['source']
            contexts.append(f"【来源：{source}】\n{doc}")
        return "\n\n---\n\n".join(contexts)
    return None
def generate_answer(query, context):
    if context is None:
        return "我在知识库中没有找到与您问题相关的内容。请尝试上传更多相关文档。"
    messages = [
        {"role": "system", "content": "你是一个专业的文档问答助手。请基于提供的文档内容回答用户的问题。如果文档内容中没有相关信息，请明确告知用户。"},
        {"role": "user", "content": f"【文档参考内容】\n{context}\n\n【用户问题】\n{query}\n\n请基于上述文档内容回答问题："}
    ]
    response = client.chat.completions.create(
        model="glm-4-flash",
        messages=messages
    )
    return response.choices[0].message.content
st.set_page_config(page_title="智能文档问答系统", layout="wide")
st.title("📚 智能文档问答系统")
st.caption("支持 PDF、Word、PPT、Excel、TXT 格式，AI 将基于文档内容回答你的问题")
chroma_client, collection, embedding_fn = init_chroma()
with st.sidebar:
    st.header("📤 上传文档")
    uploaded_files = st.file_uploader(
        "选择文件（PDF / Word / PPT / Excel / TXT）",
        type=['pdf', 'docx', 'pptx', 'xlsx', 'txt'],
        accept_multiple_files=True
    )
    if uploaded_files:
        for file in uploaded_files:
            with st.spinner(f"正在处理 {file.name}..."):
                try:
                    chunk_count, msg = upload_document(file, collection, embedding_fn)
                    if chunk_count > 0:
                        st.success(f"✅ {file.name} 已上传（{chunk_count} 个片段）")
                    else:
                        st.error(f"❌ {file.name} 处理失败：{msg}")
                except Exception as e:
                    st.error(f"❌ {file.name} 处理失败：{str(e)}")

    st.divider()
    if st.button("🗑️ 清空所有文档"):
        try:
            chroma_client.delete_collection("documents")
            st.rerun()
        except:
            st.warning("知识库已清空")
st.header("💬 提问")
try:
    all_data = collection.get()
    if all_data and all_data.get('metadatas'):
        sources = list(set([m['source'] for m in all_data['metadatas']]))
        st.info(f"📄 当前知识库包含 {len(sources)} 份文档：{', '.join(sources)}")
    else:
        st.warning("⚠️ 知识库为空，请先上传文档")
except:
    st.warning("⚠️ 知识库为空，请先上传文档")
query = st.text_input("请输入你的问题：")
if query:
    if st.button("提交问题"):
        with st.spinner("🔍 正在检索相关文档..."):
            context = search_documents(query, collection, embedding_fn)
        if context:
            with st.spinner("🤖 AI正在生成回答..."):
                answer = generate_answer(query, context)
                st.markdown("### 📝 回答")
                st.markdown(answer)
                with st.expander("📖 查看参考文档原文"):
                    st.text(context)
        else:
            st.warning("未找到相关文档内容，请尝试上传相关文档")
